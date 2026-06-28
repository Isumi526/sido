#!/usr/bin/env python3
# 営業用 機能詳細スライド(.pptx)を生成。各機能1スライド＋スクショ差し込み。
# 実行: /tmp/pptxenv/bin/python scripts/build-sales-deck.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SHOTS = "docs/sales/screenshots"
OUT = "docs/sales/feature-deck.pptx"

GREEN = RGBColor(0x06, 0xC7, 0x55)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x55, 0x60, 0x70)
DARK  = RGBColor(0x11, 0x11, 0x11)
LIGHT = RGBColor(0xF3, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Hiragino Kaku Gothic ProN"
    return tb

def bullets(slide, x, y, w, h, items, size=15, color=INK, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r = p.add_run(); r.text = "● "; r.font.size = Pt(size); r.font.color.rgb = GREEN; r.font.bold = True
        r.font.name = "Hiragino Kaku Gothic ProN"
        r2 = p.add_run(); r2.text = head; r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.bold = True
        r2.font.name = "Hiragino Kaku Gothic ProN"
        if body:
            r3 = p.add_run(); r3.text = "  " + body; r3.font.size = Pt(size-2); r3.font.color.rgb = SUB; r3.font.bold = False
            r3.font.name = "Hiragino Kaku Gothic ProN"
    return tb

def pic_framed(slide, path, x, y, w):
    """画像を幅w固定で配置し枠を付ける。高さはアスペクト比から算出。"""
    if not os.path.exists(path):
        add_rect(slide, x, y, w, Inches(3.4), LIGHT)
        txt(slide, x, y, w, Inches(3.4), [("（画面キャプチャ）", 14, SUB, False)], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    h = int(w * ih / iw)
    frame = add_rect(slide, x - Emu(9525), y - Emu(9525), w + Emu(19050), h + Emu(19050), RGBColor(0xE3,0xE8,0xEE))
    slide.shapes.add_picture(path, x, y, width=w, height=h)

def header_bar(slide, kicker, title):
    add_rect(slide, 0, 0, SW, Inches(1.15), DARK)
    add_rect(slide, 0, 0, Inches(0.18), Inches(1.15), GREEN)
    txt(slide, Inches(0.5), Inches(0.16), Inches(11), Inches(0.3),
        [(kicker, 12, RGBColor(0x9A,0xA5,0xB1), True)])
    txt(slide, Inches(0.5), Inches(0.42), Inches(11.5), Inches(0.6),
        [(title, 26, WHITE, True)])
    # badge
    b = add_rect(slide, SW - Inches(2.0), Inches(0.36), Inches(1.55), Inches(0.42), RGBColor(0x0a,0x3a,0x22))
    txt(slide, SW - Inches(2.0), Inches(0.40), Inches(1.55), Inches(0.34),
        [("✓ 本番稼働", 12, GREEN, True)], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

def feature_slide(kicker, title, lead, items, shot):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header_bar(s, kicker, title)
    # lead line
    txt(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
        [(lead, 16, SUB, True)], sp_after=0)
    # left text
    bullets(s, Inches(0.5), Inches(2.25), Inches(5.7), Inches(4.8), items, size=15)
    # right image
    pic_framed(s, os.path.join(SHOTS, shot), Inches(6.55), Inches(2.2), Inches(6.3))
    return s

# ---- タイトル ----
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, DARK)
add_rect(s, 0, SH - Inches(0.12), SW, Inches(0.12), GREEN)
txt(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.7), [("GENLINKS", 44, GREEN, True)])
txt(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.4),
    [("内装施工の台帳を、まるごと自動化。", 30, WHITE, True),
     ("職人は LINE で日報を出すだけ。原価・見積・請求・勤怠・工程まで一気通貫。", 17, RGBColor(0xC2,0xCB,0xD6), False)], sp_after=10)
txt(s, Inches(0.9), SH - Inches(0.85), Inches(11), Inches(0.4),
    [("実装機能ガイド（営業用）｜ 実機コード監査ベース・本番稼働機能のみ掲載", 12, RGBColor(0x8A,0x95,0xA1), False)])

# ---- 全体像 ----
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header_bar(s, "OVERVIEW", "3つの価値")
cards = [
    ("📲 入力は LINE だけ・AIが読み取り", "職人はLINEで日報を出すだけ。領収書・請求書・価格表はAIが自動入力。現場の手間を最小化。"),
    ("🔗 1つのデータで台帳が全部つながる", "日報→原価→現場別集計→請求が同じ数字。二重入力なし。バラバラだった台帳を1つに。"),
    ("✉️ 取引のやり取りまで自動", "見積・発注・請求・督促をメール＆業者ポータルで。承諾は署名証跡つき。"),
]
cx = Inches(0.5)
for i,(h,b) in enumerate(cards):
    x = Inches(0.5) + i*Inches(4.2)
    add_rect(s, x, Inches(1.9), Inches(3.95), Inches(4.4), LIGHT)
    add_rect(s, x, Inches(1.9), Inches(3.95), Inches(0.12), GREEN)
    txt(s, x+Inches(0.3), Inches(2.3), Inches(3.4), Inches(1.2), [(h, 17, INK, True)], sp_after=0)
    txt(s, x+Inches(0.3), Inches(3.5), Inches(3.4), Inches(2.6), [(b, 13.5, SUB, False)], sp_after=0)

# ---- 各機能スライド ----
feature_slide("日次管理", "ダッシュボード", "月次の経営数字を、開いた瞬間に。",
    [("原価カテゴリ別集計", "社員・商社・業者・経費を月単位で"),
     ("明細ドリルダウン", "集計行クリックで日付・対象・金額まで"),
     ("通常/残業/深夜の自動再計算", "日報実績から正確に"),
     ("出張費・ガソリン代も自動集計", "")],
    "01-dashboard.png")

feature_slide("日次管理", "日報一覧", "届いた日報を一元管理。人件費は自動計算。",
    [("人件費の自動計算", "日当/時間給・残業・深夜・日曜割増を自動"),
     ("LINE通知", "未通知の日報をワンクリックで通知"),
     ("経費内訳の確認", "車両・電車・宿泊・その他まで明細で"),
     ("作業員フィルタ・月ナビ", "")],
    "08-reports.png")

feature_slide("日次管理", "現場別集計", "現場ごとの人工・経費・請求をまとめて。",
    [("現場タブで横断", "五十音順で全現場を切り替え"),
     ("CSV＋見積書PDFでエクスポート", "期間指定（当月/範囲/全期間）"),
     ("作業員別・経費別の明細ドロワー", ""),
     ("月計・年計の集計", "")],
    "07-site-reports.png")

feature_slide("工程", "工程管理", "複数現場の工程を、エクセル風ガントで統合。",
    [("全現場横断ビュー", "カレンダー軸で納期を可視化"),
     ("並び替え", "現場別/担当別/開始日順"),
     ("工種で色分け", "日中/夜間/家具などを一目で"),
     ("契約金額・現場管理者・メモ", "1現場に複数工程(1:n)も")],
    "02-process.png")

feature_slide("見積・発注", "見積もり（作成〜発注）", "明細を組んで見積書、そのまま商社へ発注。",
    [("見積ビルダー", "場所・工種・品名・単価で明細を作成→PDF発行"),
     ("受注で現場化", "見積→現場へ昇華"),
     ("商社へ発注書を自動生成・メール送信", ""),
     ("マスタ連携", "材料・工種・単価・法定福利費・消費税")],
    "03-estimate-list.png")

feature_slide("見積・発注", "注文書発行・承諾", "業者がリンクで承諾。署名証跡まで残る。",
    [("注文書をメール送信", "業者は承諾用リンクで署名・同意"),
     ("変更注文書の再承諾", "金額増減を再依頼"),
     ("請求依頼まで一気通貫", "承諾済みに請求フォームURLを送付"),
     ("承諾の証跡", "署名画像・IP・PDFハッシュを保存")],
    "05-purchase-orders.png")

feature_slide("経費・請求", "協力業者請求", "下請けの請求を受領・管理。AIで自動入力。",
    [("請求書PDFをAI解析", "業者名・請求日・金額・明細を自動抽出"),
     ("注文書残額と照合", "超過を自動で弾く"),
     ("支払期限・支払状況の管理", "未払い/支払い済みタブ"),
     ("明細ごとの領収書添付", "")],
    "06-subcontractor-invoices.png")

feature_slide("経費・原価", "ガソリン按分", "燃料代を、現場の走行距離比で配賦。",
    [("月次実費を自動集計", "日報の「本日のガソリン代」から"),
     ("走行距離も自動集計", "日報の車両経費から"),
     ("見込み vs 実績 vs 差異", "並べて確認"),
     ("手動上書き", "集計が実態と合わない時")],
    "09-gasoline.png")

feature_slide("勤怠・労務", "出面・勤怠", "勤務時間と人件費を、給与計算用に集計。",
    [("時間帯別の集計", "通常・残業・深夜・日曜・休日を正確に"),
     ("サマリーカード", "稼働日数・休み・出張日数"),
     ("人件費の表示/非表示トグル", ""),
     ("PDF出力（印刷版）", "")],
    "10-worker-reports.png")

feature_slide("マスタ管理", "作業員マスタ", "職員情報・賃金・履歴を一元管理。",
    [("昇給履歴", "発効日付きで過去日報も正しく計算"),
     ("雇用形態・週所定日数", "正社員/パート/業務委託"),
     ("家族構成・車検・健診履歴", ""),
     ("ログイン認証（email/pass）", "作業員ごとに設定")],
    "11-workers.png")

feature_slide("勤怠・労務", "有給管理", "法定5日取得義務まで、まとめて管理。",
    [("付与・使用・残日数", "作業員別サマリー"),
     ("法定5日義務の達成度", "コンプライアンス対応"),
     ("付与歴（発効日・理由・日数）", ""),
     ("管理簿の印刷", "")],
    "12-paid-leave.png")

# ---- LIFF（スクショ別途）----
s = feature_slide("作業員アプリ", "LINE 日報アプリ（LIFF）", "アプリ不要。LINEだけで現場が完結。",
    [("LINE日報入力", "稼働区分→現場→経費→送信。複数現場を1日でまとめて"),
     ("領収書AI自動入力", "撮るだけで支払先・金額・インボイス番号を抽出"),
     ("経費の複数明細・個人立替", "宿泊/駐車/高速/電車を明細ごと"),
     ("下書き自動保存", "中断しても画像ごと復元"),
     ("出退勤チェックイン(QR)", "安全ルール同意・GPS証跡"),
     ("代理入力", "LINEを持たない職人の分も")],
    "_liff_placeholder.png")

# ---- 裏側の自動化（テキスト主体）----
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header_bar(s, "AUTOMATION", "裏側で自動で動く仕組み")
txt(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
    [("人手をかけずに、通知・解析・送信が回り続ける。", 16, SUB, True)], sp_after=0)
auto = [
    ("日報未送信リマインド", "設定時刻に未送信者を検知しLINE個人通知（会社単位で時刻/対象/ON-OFF）"),
    ("車検期日リマインド", "車検が近い車両を週次で検知して通知"),
    ("領収書・請求書・価格表のAI解析", "撮影/PDFから自動抽出。価格表は単価マスタへ（承認制）"),
    ("各種メール自動送信", "見積・注文書・変更注文・請求依頼・業者登録・経費申請PDF"),
    ("下請けポータル", "業者がログイン不要のリンクで承諾・請求・見積アップロード（署名証跡）"),
    ("マルチテナント運用基盤", "会社（テナント）単位で通知・設定を分離"),
]
bullets(s, Inches(0.6), Inches(2.3), Inches(12), Inches(4.6), auto, size=16, gap=12)

# ---- クロージング ----
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, DARK)
add_rect(s, 0, 0, Inches(0.18), SH, GREEN)
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.6), [("GENLINKS", 30, GREEN, True)])
txt(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(2.2),
    [("掲載した機能は、すべて実機のソースコードで実装・配線を確認済み。", 18, WHITE, True),
     ("admin 30画面・LINE日報フロー全体・バックエンド23本を監査。", 15, RGBColor(0xC2,0xCB,0xD6), False),
     ("構想中・未実装の項目は載せていません（盛らない方針）。", 15, RGBColor(0xC2,0xCB,0xD6), False)], sp_after=10)
txt(s, Inches(0.9), SH - Inches(0.9), Inches(11), Inches(0.4),
    [("2026-06-28 時点 ｜ 画面キャプチャはローカル環境（デモデータは適宜差し替え）", 11, RGBColor(0x8A,0x95,0xA1), False)])

prs.save(OUT)
print("saved", OUT, "/ slides:", len(prs.slides._sldIdLst))
